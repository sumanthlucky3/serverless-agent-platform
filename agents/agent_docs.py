"""
agents/agent_docs.py
=====================
Docs & Reports Agent — generates professional documents from a prompt.

Capabilities (v1):
  - Word documents (.docx)     via python-docx
  - Excel spreadsheets (.xlsx) via openpyxl
  - PowerPoint slides (.pptx)  via python-pptx
  - PDF reports                via reportlab
  - Plain Markdown (.md)       built-in

Powered by: Google Gemini (content generation) +
            python-docx / openpyxl / python-pptx (file creation)

Future upgrades:
  - Upload generated file to Supabase Storage
  - Return download link in output
  - Support image generation (charts, diagrams)
"""

import os
import json
import re
from google.antigravity import Agent, LocalAgentConfig


# ── Helpers — file generators ────────────────────────────────────────────────

def _write_docx(content: str, filename: str) -> str:
    """Write a Word .docx file from markdown-style content."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("# "):
            p = doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith("## "):
            p = doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("### "):
            p = doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped == "":
            doc.add_paragraph()
        else:
            doc.add_paragraph(stripped)
    doc.save(filename)
    return filename


def _write_xlsx(content: str, filename: str) -> str:
    """Write an Excel .xlsx file. Expects content with pipe-separated rows."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Agent Report"

    header_fill = PatternFill(start_color="1E40AF", end_color="1E40AF", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")

    lines = [l for l in content.split("\n") if "|" in l and l.strip()]
    for ri, line in enumerate(lines):
        cells = [c.strip() for c in line.split("|") if c.strip()]
        for ci, val in enumerate(cells):
            cell = ws.cell(row=ri+1, column=ci+1, value=val)
            if ri == 0:
                cell.fill = header_fill
                cell.font = header_font
            cell.alignment = Alignment(wrap_text=True)

    # Auto-width columns
    for col in ws.columns:
        max_len = max((len(str(c.value or "")) for c in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    wb.save(filename)
    return filename


def _write_pptx(content: str, filename: str) -> str:
    """Write a PowerPoint .pptx from content with slide markers."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    prs = Presentation()
    slide_layout_title = prs.slide_layouts[0]  # Title slide
    slide_layout_body  = prs.slide_layouts[1]  # Title + content

    slides_raw = re.split(r"(?m)^---+$", content)
    for i, slide_text in enumerate(slides_raw):
        lines = [l.strip() for l in slide_text.strip().split("\n") if l.strip()]
        if not lines:
            continue
        title = lines[0].lstrip("#").strip()
        body  = "\n".join(lines[1:])

        layout = slide_layout_title if i == 0 else slide_layout_body
        slide  = prs.slides.add_slide(layout)

        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1 and body:
            slide.placeholders[1].text = body

    prs.save(filename)
    return filename


def _write_pdf(content: str, filename: str) -> str:
    """Write a PDF report from plain text content."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.lib.enums import TA_LEFT

    doc    = SimpleDocTemplate(filename, pagesize=A4,
                               leftMargin=2.5*cm, rightMargin=2.5*cm,
                               topMargin=2*cm, bottomMargin=2*cm)
    styles = getSampleStyleSheet()
    story  = []

    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("# "):
            story.append(Paragraph(stripped[2:], styles["Heading1"]))
            story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#16A34A")))
        elif stripped.startswith("## "):
            story.append(Paragraph(stripped[3:], styles["Heading2"]))
        elif stripped.startswith("### "):
            story.append(Paragraph(stripped[4:], styles["Heading3"]))
        elif stripped == "":
            story.append(Spacer(1, 0.3*cm))
        else:
            story.append(Paragraph(stripped, styles["Normal"]))

    doc.build(story)
    return filename


# ── Main agent entry point ────────────────────────────────────────────────────

async def run(task: str, supabase, agent_id: str, session_id=None) -> dict:
    print("[AGENT_DOCS] Starting Docs & Reports Agent...")

    gemini_key = os.getenv("GEMINI_API_KEY")
    if not gemini_key:
        raise ValueError("GEMINI_API_KEY is required for agent_docs")

    # ── Detect desired output format from task ────────────────────
    task_lower = task.lower()
    if any(w in task_lower for w in ["excel", "spreadsheet", "xlsx", "table", "csv"]):
        output_format = "xlsx"
        format_hint   = "pipe-separated table rows (first row = headers)"
    elif any(w in task_lower for w in ["powerpoint", "ppt", "pptx", "presentation", "slides"]):
        output_format = "pptx"
        format_hint   = "slides separated by '---', each starting with a title line"
    elif any(w in task_lower for w in ["pdf"]):
        output_format = "pdf"
        format_hint   = "well-structured markdown text with # headings"
    else:
        output_format = "docx"
        format_hint   = "well-structured markdown text with # headings and bullet points"

    print(f"[AGENT_DOCS] Detected output format: {output_format}")

    # ── Step 1: Ask Gemini to generate structured content ─────────
    system_context = f"""You are a professional document writer.
Generate content for a {output_format.upper()} document based on the user's request.
Format your output as {format_hint}.
Be thorough, professional, and well-structured. Include relevant sections, details, and examples.
Output ONLY the document content — no explanations, no preamble."""

    config = LocalAgentConfig(api_key=gemini_key)
    async with Agent(config) as agent:
        response    = await agent.chat(f"{system_context}\n\nTask: {task}")
        doc_content = await response.text()

    print(f"[AGENT_DOCS] Gemini generated {len(doc_content)} chars of content.")

    # ── Step 2: Create the file ───────────────────────────────────
    safe_name   = re.sub(r"[^a-zA-Z0-9_]", "_", task[:40])
    output_file = f"output_{safe_name}.{output_format}"

    try:
        if output_format == "docx":
            _write_docx(doc_content, output_file)
        elif output_format == "xlsx":
            _write_xlsx(doc_content, output_file)
        elif output_format == "pptx":
            _write_pptx(doc_content, output_file)
        elif output_format == "pdf":
            _write_pdf(doc_content, output_file)
        print(f"[AGENT_DOCS] File created: {output_file}")
    except Exception as e:
        print(f"[AGENT_DOCS] File creation failed: {e}. Falling back to plain text.")
        output_file = "output.txt"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(doc_content)

    # ── Step 3: Log to Supabase ───────────────────────────────────
    summary_text = f"[{output_format.upper()} generated] {output_file}\n\n{doc_content[:500]}..."
    if supabase:
        try:
            supabase.table("agent_runs").insert({
                "task":       task,
                "routed_to":  "DOCS",
                "result":     summary_text,
                "model_used": "gemini/antigravity-preview-05-2026",
                "tokens_used": 0,
                "status":     "success"
            }).execute()

            # Log file record
            supabase.table("agent_files").insert({
                "session_id": int(session_id) if session_id else None,
                "agent_id":   agent_id,
                "file_name":  output_file,
                "file_type":  output_format,
                "direction":  "output"
            }).execute()

            print("[AGENT_DOCS] Logged to Supabase.")
        except Exception as e:
            print(f"[AGENT_DOCS] Warning: Supabase log failed — {e}")

    return {
        "output_text":  summary_text,
        "output_file":  output_file,
        "output_format": output_format,
        "model_used":   "gemini/antigravity-preview-05-2026",
        "status":       "success"
    }
