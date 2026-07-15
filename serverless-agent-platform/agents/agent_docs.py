"""
agents/agent_docs.py
=====================
Docs & Reports Agent — generates professional documents from a prompt.

Powered by: NVIDIA Nemotron via OpenRouter API (content generation) +
            python-docx / openpyxl / python-pptx (file creation)
"""

import os
import re
import httpx


# ── Helpers — file generators ────────────────────────────────────────────────

def _write_docx(content: str, filename: str) -> str:
    """Write a Word .docx file from markdown-style content."""
    try:
        from docx import Document
        doc = Document()
        for line in content.split("\n"):
            stripped = line.strip()
            if stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped == "":
                doc.add_paragraph()
            else:
                doc.add_paragraph(stripped)
        doc.save(filename)
    except ImportError:
        print("[AGENT_DOCS] python-docx not installed. Falling back to plain text.")
        raise
    return filename


def _write_xlsx(content: str, filename: str) -> str:
    """Write an Excel .xlsx file from CSV-style content."""
    try:
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        for row in content.split("\n"):
            if row.strip():
                ws.append(row.split(","))
        wb.save(filename)
    except ImportError:
        print("[AGENT_DOCS] openpyxl not installed. Falling back to plain text.")
        raise
    return filename


def _write_pptx(content: str, filename: str) -> str:
    """Write a PowerPoint .pptx file from markdown-style content."""
    try:
        from pptx import Presentation
        prs = Presentation()
        slides = content.split("---")
        for slide_text in slides:
            if not slide_text.strip():
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines = [l.strip() for l in slide_text.strip().split("\n") if l.strip()]
            if lines:
                slide.shapes.title.text = lines[0].replace("#", "").strip()
                tf = slide.placeholders[1].text_frame
                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.text = line.replace("-", "").replace("*", "").strip()
        prs.save(filename)
    except ImportError:
        print("[AGENT_DOCS] python-pptx not installed. Falling back to plain text.")
        raise
    return filename


def _write_pdf(content: str, filename: str) -> str:
    """Write a PDF file."""
    try:
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(filename)
        y = 800
        for line in content.split("\n"):
            c.drawString(50, y, line.strip()[:100])
            y -= 20
            if y < 50:
                c.showPage()
                y = 800
        c.save()
    except ImportError:
        print("[AGENT_DOCS] reportlab not installed. Falling back to plain text.")
        raise
    return filename


async def run(task: str, supabase, agent_id: str, session_id=None) -> dict:
    print("[AGENT_DOCS] Starting Docs & Reports Agent...")

    openrouter_key = os.getenv("OPENROUTER_API_KEY") or os.getenv("VITE_OPENROUTER_API_KEY")
    if not openrouter_key:
        raise ValueError("OPENROUTER_API_KEY is required for agent_docs")

    # ── Step 0: Determine output format from task ─────────────────
    lower_task = task.lower()
    output_format = "md"
    format_hint = "plain markdown"
    
    if "word" in lower_task or "docx" in lower_task or "document" in lower_task:
        output_format = "docx"
        format_hint = "structured text with # headings and bullet points"
    elif "excel" in lower_task or "xlsx" in lower_task or "spreadsheet" in lower_task or "csv" in lower_task:
        output_format = "xlsx"
        format_hint = "comma-separated values (CSV) format"
    elif "powerpoint" in lower_task or "pptx" in lower_task or "presentation" in lower_task or "slide" in lower_task:
        output_format = "pptx"
        format_hint = "slides separated by '---' with a # Title and bullet points"
    elif "pdf" in lower_task:
        output_format = "pdf"
        format_hint = "plain text suitable for a PDF report"

    print(f"[AGENT_DOCS] Detected format: {output_format}")

    # ── Step 1: Ask LLM to generate structured content ─────────
    system_context = f"""You are a professional document writer.
Generate content for a {output_format.upper()} document based on the user's request.
Format your output as {format_hint}.
Be thorough, professional, and well-structured. Include relevant sections, details, and examples.
Output ONLY the document content — no explanations, no preamble."""

    hf_key = os.getenv("HUGGINGFACE_API_KEY") or os.getenv("VITE_HUGGINGFACE_API_KEY")
    openrouter_key = os.getenv("OPENROUTER_API_KEY") or os.getenv("VITE_OPENROUTER_API_KEY")
    routed_model = os.getenv("HF_ROUTED_MODEL_ID")

    doc_content = ""
    model_used = ""

    # ── Call Hugging Face API ────────────────────────────────────
    if hf_key and routed_model:
        print(f"[AGENT_DOCS] Using Hugging Face model: {routed_model}")
        headers = {"Authorization": f"Bearer {hf_key}", "Content-Type": "application/json"}
        payload = {
            "inputs": f"<|system|>\n{system_context}<|end|>\n<|user|>\n{task}<|end|>\n<|assistant|>\n",
            "parameters": {"max_new_tokens": 2048, "temperature": 0.2, "return_full_text": False}
        }
        
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(
                    f"https://api-inference.huggingface.co/models/{routed_model}",
                    headers=headers,
                    json=payload
                )
                response.raise_for_status()
                data = response.json()
                if isinstance(data, list) and "generated_text" in data[0]:
                    doc_content = data[0]["generated_text"].strip()
                elif "error" in data:
                    raise Exception(data["error"])
                else:
                    doc_content = str(data)
                model_used = routed_model
        except Exception as e:
            print(f"[AGENT_DOCS] HF Inference failed: {e}. Falling back to OpenRouter...")
            hf_key = None # Force fallback

    # ── Fallback to OpenRouter ───────────────────────────────────
    if not doc_content and openrouter_key:
        print("[AGENT_DOCS] Using OpenRouter fallback (NVIDIA Nemotron)...")
        headers = {
            "Authorization": f"Bearer {openrouter_key}",
            "Content-Type": "application/json",
            "HTTP-Referer": "https://sumanthlucky3.github.io/serverless-agent-platform/",
            "X-Title": "Serverless Agent Platform",
        }
        payload = {
            "model": "nvidia/llama-3.1-nemotron-ultra-253b-v1:free",
            "messages": [
                {"role": "system", "content": system_context},
                {"role": "user", "content": task}
            ],
            "max_tokens": 2048,
        }

        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post("https://openrouter.ai/api/v1/chat/completions", headers=headers, json=payload)
            response.raise_for_status()
            data = response.json()

        doc_content = data["choices"][0]["message"]["content"]
        model_used = data.get("model", "nvidia/nemotron-ultra")

    if not doc_content:
        raise ValueError("Neither Hugging Face nor OpenRouter keys are available, or both APIs failed.")

    print(f"[AGENT_DOCS] Content generated ({len(doc_content)} chars).")

    # ── Step 2: Create the file ───────────────────────────────────
    safe_name = re.sub(r"[^a-zA-Z0-9_]", "_", task[:20])
    if not safe_name: safe_name = "report"
    output_file = f"output_{safe_name}.{output_format}"

    try:
        if output_format == "docx":
            _write_docx(doc_content, output_file)
        elif output_format == "xlsx":
            _write_xlsx(doc_content, output_file)
        elif output_format == "pptx":
            _write_pptx(doc_content, output_file)
        elif output_format == "pdf":
            _write_pdf(doc_content, output_file)
        elif output_format == "md":
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(doc_content)
        print(f"[AGENT_DOCS] File created: {output_file}")
    except Exception as e:
        print(f"[AGENT_DOCS] File creation failed: {e}. Falling back to plain text.")
        output_file = "output.txt"
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(doc_content)

    # ── Step 3: Log file to Supabase ───────────────────────────────────
    summary_text = f"[{output_format.upper()} generated] {output_file}\n\n{doc_content[:500]}..."
    if supabase:
        try:
            # Log file record
            supabase.table("agent_files").insert({
                "session_id": int(session_id) if session_id else None,
                "agent_id":   agent_id,
                "file_name":  output_file,
                "file_type":  output_format,
                "direction":  "output"
            }).execute()

            print("[AGENT_DOCS] File logged to Supabase.")
        except Exception as e:
            print(f"[AGENT_DOCS] Warning: Supabase log failed — {e}")

    return {
        "output_text":  summary_text,
        "output_file":  output_file,
        "output_format": output_format,
        "model_used":   model_used,
        "status":       "success"
    }
